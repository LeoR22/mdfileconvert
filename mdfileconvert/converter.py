import pdfplumber
import re
import io
import os
from PIL import Image
import docx
import pandas as pd
import json
import xml.etree.ElementTree as ET
from pptx import Presentation

# Umbrales mínimos (en puntos) para considerar una imagen "grande y entendible"
MIN_WIDTH = 100
MIN_HEIGHT = 100

def clean_cell(cell):
    """
    Limpia el contenido de la celda:
      - Si es None, devuelve una cadena vacía.
      - Une líneas separadas en una sola línea.
      - Elimina espacios excesivos.
    """
    if cell is None:
        return ""
    return " ".join(cell.splitlines()).strip()

def auto_format_text(text):
    """
    Aplica heurísticas sencillas para transformar el texto extraído en Markdown:
      - Si una línea es corta (menos de 60 caracteres) y está en mayúsculas,
        se asume que es un título y se le antepone "## ".
      - Si una línea es corta y está en formato Title Case, se le antepone "### ".
      - Si la línea empieza con un guión (listas), se deja sin cambios.
      - El resto se trata como párrafo.
    """
    lines = text.splitlines()
    new_lines = []
    
    for line in lines:
        stripped = line.strip()
        if not stripped:
            new_lines.append("")  # Conservar línea en blanco para separar párrafos
            continue
        
        # Si la línea parece ser un título (corta y todo en mayúsculas)
        if len(stripped) < 60 and stripped == stripped.upper():
            new_lines.append("## " + stripped)
            continue
        
        # Si la línea parece ser un subtítulo (corta y en Title Case)
        if len(stripped) < 60 and stripped.istitle():
            new_lines.append("### " + stripped)
            continue
        
        # Si la línea empieza con indicadores de lista, se deja como está
        if re.match(r"^[\-\*\+]\s+", stripped):
            new_lines.append(stripped)
        else:
            new_lines.append(stripped)
    
    # Unir líneas con doble salto de línea para separar párrafos
    return "\n\n".join(new_lines)

def pdf_to_markdown(pdf_file):
    """
    Convierte un archivo PDF (objeto file-like) a Markdown.
    Extrae texto, aplica formato, procesa tablas e imágenes.
    Las imágenes se guardan en la carpeta "temporal" si son lo suficientemente grandes.
    """
    output = io.StringIO()
    
    # Crear la carpeta temporal para las imágenes si no existe
    os.makedirs("temporal", exist_ok=True)
    
    with pdfplumber.open(pdf_file) as pdf:
        for page_num, page in enumerate(pdf.pages, start=1):
            output.write(f"## Página {page_num}\n\n")
            
            # Extraer y formatear el texto
            text = page.extract_text()
            if text:
                formatted_text = auto_format_text(text)
                output.write(formatted_text + "\n\n")
            
            # Extraer y formatear las tablas
            tables = page.extract_tables()
            if tables:
                output.write("### Tablas extraídas:\n\n")
                for table in tables:
                    if not table:
                        continue
                    
                    # Calcular el número máximo de columnas
                    max_cols = max(len(row) for row in table if row is not None)
                    
                    # Normalizar la tabla: limpiar celdas y rellenar filas incompletas
                    normalized_table = []
                    for row in table:
                        if row:
                            new_row = [clean_cell(cell) for cell in row]
                            if len(new_row) < max_cols:
                                new_row.extend([""] * (max_cols - len(new_row)))
                            normalized_table.append(new_row)
                    
                    # Formatear la tabla en sintaxis Markdown (asumiendo que la primera fila es cabecera)
                    if normalized_table:
                        header = normalized_table[0]
                        output.write("| " + " | ".join(header) + " |\n")
                        output.write("|" + "|".join([":-:" for _ in range(max_cols)]) + "|\n")
                        for row in normalized_table[1:]:
                            output.write("| " + " | ".join(row) + " |\n")
                        output.write("\n")
            
            # Extraer imágenes de la página, si las hay
            if hasattr(page, "images") and page.images:
                for j, image in enumerate(page.images):
                    # Calcular ancho y alto de la imagen
                    width = image["x1"] - image["x0"]
                    height = image["bottom"] - image["top"]
                    
                    # Solo procesar la imagen si cumple con los umbrales mínimos
                    if width < MIN_WIDTH or height < MIN_HEIGHT:
                        continue
                    
                    # Definir el recorte según las coordenadas de la imagen
                    bbox = (image["x0"], image["top"], image["x1"], image["bottom"])
                    
                    # Obtener la imagen de la página
                    page_img = page.to_image()
                    # Utilizar la imagen original (objeto PIL) para recortar
                    cropped = page_img.original.crop(bbox)
                    
                    # Definir un nombre único para la imagen
                    image_filename = f"temporal/imagen_{page_num}_{j}.png"
                    cropped.save(image_filename, format="PNG")
                    
                    # Añadir la referencia a la imagen en Markdown
                    output.write(f"![Imagen de página {page_num}, imagen {j}]({image_filename})\n\n")
                    
    return output.getvalue()

def docx_to_markdown(docx_file):
    """Convierte un archivo Word (.docx) a formato Markdown."""
    doc = docx.Document(docx_file)
    output = []
    
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            if para.style.name.startswith("Heading"):
                level = int(para.style.name[-1]) if para.style.name[-1].isdigit() else 2
                output.append("#" * level + " " + text)
            else:
                output.append(text)
    
    return "\n\n".join(output)



def xlsx_to_markdown(xlsx_file):
    """Convierte un archivo Excel (.xlsx) a formato Markdown."""
    output = []
    xls = pd.ExcelFile(xlsx_file)
    
    for sheet_name in xls.sheet_names:
        df = xls.parse(sheet_name)
        output.append(f"## Hoja: {sheet_name}\n")
        output.append(df.to_markdown(index=False))
    
    return "\n\n".join(output)

def pptx_to_markdown(pptx_file):
    """Convierte un archivo PowerPoint (.pptx) a formato Markdown."""
    prs = Presentation(pptx_file)
    output = []
    
    for i, slide in enumerate(prs.slides, start=1):
        output.append(f"## Diapositiva {i}\n")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    output.append(text)
    
    return "\n\n".join(output)

def csv_to_markdown(csv_file):
    """Convierte un archivo CSV a formato Markdown."""
    df = pd.read_csv(csv_file)
    return df.to_markdown(index=False)

def json_to_markdown(json_file):
    """Convierte un archivo JSON a formato Markdown como bloque de código."""
    data = json.load(json_file)
    return "```json\n" + json.dumps(data, indent=4) + "\n```"

def xml_to_markdown(xml_file):
    """Convierte un archivo XML a formato Markdown."""
    tree = ET.parse(xml_file)
    root = tree.getroot()
    
    def parse_element(element, level=0):
        md = "#" * (level + 2) + f" {element.tag}\n"
        if element.text and element.text.strip():
            md += f"{element.text.strip()}\n"
        for child in element:
            md += parse_element(child, level + 1)
        return md
    return parse_element(root)